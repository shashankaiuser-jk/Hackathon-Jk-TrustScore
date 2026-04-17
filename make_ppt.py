"""
CreditLens AI — Demo Deck Generator
Run: python3 make_ppt.py
Output: CreditLens_AI_Demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x16, 0xA3, 0x4A)   # brand green
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_DARK  = RGBColor(0x0F, 0x1D, 0x0F)   # sidebar dark
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF8, 0xFA, 0xF8)
BORDER      = RGBColor(0xE3, 0xED, 0xE3)
MUTED       = RGBColor(0x6B, 0x7A, 0x6B)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
GOLD        = RGBColor(0xF5, 0x9E, 0x0B)
TEXT_DARK   = RGBColor(0x11, 0x1B, 0x11)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color: RGBColor):
    """Fill slide background."""
    bg_elem = slide.background
    fill    = bg_elem.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0), radius=None):
    """Add a rectangle / rounded-rect shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()           # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = line_width
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    run.font.name      = "Calibri"
    return txb


def mtxt(slide, lines, x, y, w, h,
         size=16, color=TEXT_DARK, align=PP_ALIGN.LEFT, line_spacing=1.15):
    """Multi-line textbox from list of (text, bold, size_override, color_override)."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, False, None, None)
        text, bold, sz, col = item[0], item[1] if len(item)>1 else False, item[2] if len(item)>2 else None, item[3] if len(item)>3 else None
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(sz if sz else size)
        run.font.bold  = bold
        run.font.color.rgb = col if col else color
        run.font.name  = "Calibri"
    return txb


def pill(slide, text, x, y, w, h, fill, text_color, size=11, bold=True):
    """Badge / pill shape."""
    s = rect(slide, x, y, w, h, fill_color=fill)
    tf = s.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = text_color
    run.font.name  = "Calibri"
    # vertical centre
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    tf.word_wrap = False
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s


def divider(slide, y, color=BORDER):
    rect(slide, Inches(0.6), y, Inches(12.13), Pt(1), fill_color=color)


def section_label(slide, text, x, y):
    txt(slide, text, x, y, Inches(4), Inches(0.3),
        size=8, bold=True, color=GREEN,
        align=PP_ALIGN.LEFT)


def card(slide, x, y, w, h, fill=NEAR_WHITE, border=BORDER):
    return rect(slide, x, y, w, h, fill_color=fill, line_color=border, line_width=Pt(0.75))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GREEN_DARK)

# Left green accent bar
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=GREEN)

# THIN FILE badge top-right
pill(sl, "THIN FILE ANALYSIS MODE ACTIVE",
     Inches(9.8), Inches(0.35), Inches(3.2), Inches(0.38),
     RGBColor(0x16,0xA3,0x4A), WHITE, size=9)

# Main title
txt(sl, "CreditLens AI",
    Inches(0.7), Inches(1.8), Inches(9), Inches(1.2),
    size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Green underline accent
rect(sl, Inches(0.7), Inches(3.1), Inches(3.5), Pt(5), fill_color=GREEN)

# Subtitle
txt(sl, "Credit Intelligence for the Thin File Economy",
    Inches(0.7), Inches(3.3), Inches(9), Inches(0.7),
    size=24, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.LEFT)

# One-liner
txt(sl, "AI-powered risk scoring without bureau dependency  ·  35+ live workflows  ·  Butterfly effect detection",
    Inches(0.7), Inches(4.1), Inches(11), Inches(0.5),
    size=14, color=RGBColor(0x6B,0x7A,0x6B), align=PP_ALIGN.LEFT)

# Bottom stat strip
for i, (val, lbl) in enumerate([
    ("₹2.4T", "MSME Credit Gap"),
    ("60%",   "Borrowers w/o bureau data"),
    ("35+",   "Risk Workflows"),
    ("Zero",  "CIBIL Dependency"),
]):
    bx = Inches(0.7 + i * 3.15)
    rect(sl, bx, Inches(5.5), Inches(2.9), Inches(1.55),
         fill_color=RGBColor(0x1A,0x2E,0x1A),
         line_color=RGBColor(0x22,0x4A,0x22), line_width=Pt(0.75))
    txt(sl, val, bx + Inches(0.18), Inches(5.62), Inches(2.6), Inches(0.55),
        size=28, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    txt(sl, lbl, bx + Inches(0.18), Inches(6.15), Inches(2.6), Inches(0.4),
        size=11, color=RGBColor(0x6B,0x7A,0x6B), align=PP_ALIGN.LEFT)

# Hackathon tag bottom right
txt(sl, "Hackathon Demo  ·  2026",
    Inches(10.5), Inches(7.1), Inches(2.6), Inches(0.3),
    size=9, color=RGBColor(0x3B,0x4B,0x3B), align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=RED)

section_label(sl, "THE PROBLEM", Inches(0.55), Inches(0.3))
txt(sl, "MSME lenders are flying blind",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.9),
    size=38, bold=True, color=TEXT_DARK)

divider(sl, Inches(1.6))

# Left — problem narrative
mtxt(sl, [
    ("60% of MSMEs are thin-file borrowers", True, 15, TEXT_DARK),
    ("They have no CIBIL history, no formal credit trail.", False, 13, MUTED),
    ("", False, 6, None),
    ("Traditional credit scoring fails them.", True, 15, RED),
    ("Lenders either reject them outright or approve blindly —\nboth outcomes cost money.", False, 13, MUTED),
    ("", False, 6, None),
    ("Channel partners act as intermediaries but carry\nhidden risk that nobody is measuring.", True, 15, TEXT_DARK),
    ("One partner's default can cascade across the entire\nportfolio through butterfly effects.", False, 13, MUTED),
], Inches(0.55), Inches(1.75), Inches(5.8), Inches(4.5), size=13)

# Right — stat cards
for i, (val, lbl, sub, col) in enumerate([
    ("₹2.4T",  "MSME Credit Gap",         "Unmet formal lending demand in India", RED),
    ("60%",    "No bureau data",           "Borrowers invisible to traditional scoring", ORANGE),
    ("3–5×",   "Higher default detection", "When behavioral signals are used vs bureau", GREEN),
    ("48 hrs", "Avg manual review time",   "Replaced by CreditLens in under 5 seconds", BLUE),
]):
    cy = Inches(1.75 + i * 1.35)
    card(sl, Inches(6.7), cy, Inches(6.3), Inches(1.2), fill=NEAR_WHITE)
    rect(sl, Inches(6.7), cy, Inches(0.18), Inches(1.2), fill_color=col)
    txt(sl, val,  Inches(7.05), cy + Inches(0.1),  Inches(3), Inches(0.55), size=26, bold=True, color=col)
    txt(sl, lbl,  Inches(7.05), cy + Inches(0.6),  Inches(5.8), Inches(0.3), size=12, bold=True,  color=TEXT_DARK)
    txt(sl, sub,  Inches(7.05), cy + Inches(0.88), Inches(5.8), Inches(0.28), size=10, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=GREEN)

section_label(sl, "THE SOLUTION", Inches(0.55), Inches(0.3))
txt(sl, "CreditLens AI — What we built",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.9),
    size=36, bold=True, color=TEXT_DARK)

divider(sl, Inches(1.6))

# 6 solution pillars in 2×3 grid
pillars = [
    ("🏦", "No Bureau Dependency",    "Pure behavioral + transactional scoring.\nDigitap signals, GST, repayment history."),
    ("⚡", "35+ Risk Workflows",      "Automated rules covering repayment,\nfinancial health, behavioral & contextual risk."),
    ("🌊", "Butterfly Effect Engine", "Detects how one small event cascades\ninto portfolio-wide default risk."),
    ("🔍", "Pattern Detection",       "Surfaces hidden patterns: partial payments,\ncash flow stress, multi-loan overlap."),
    ("🤖", "AI Reasoning Layer",      "Every score comes with plain-language\nexplanation of how AI reached the verdict."),
    ("📊", "Live Health Checks",      "Run on-demand risk scoring on any loan\nor channel partner in under 5 seconds."),
]
for i, (icon, title, desc) in enumerate(pillars):
    col = i % 3
    row = i // 3
    cx = Inches(0.45 + col * 4.25)
    cy = Inches(1.85 + row * 2.55)
    card(sl, cx, cy, Inches(4.0), Inches(2.35))
    rect(sl, cx, cy, Inches(0.18), Inches(2.35), fill_color=GREEN)
    txt(sl, icon,  cx + Inches(0.35), cy + Inches(0.18), Inches(0.6),  Inches(0.5),  size=22)
    txt(sl, title, cx + Inches(0.35), cy + Inches(0.65), Inches(3.55), Inches(0.45), size=14, bold=True, color=TEXT_DARK)
    txt(sl, desc,  cx + Inches(0.35), cy + Inches(1.08), Inches(3.55), Inches(1.1),  size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS (pipeline)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GREEN_DARK)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=GREEN)

section_label(sl, "HOW IT WORKS", Inches(0.55), Inches(0.3))
txt(sl, "The AI Scoring Pipeline",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=WHITE)

# Pipeline steps
steps = [
    ("1", "Data Ingestion",      "Digitap bank signals\nGST revenue trends\nRepayment history",       BLUE),
    ("2", "Signal Processing",   "Transaction pattern analysis\nCash flow mapping\nAnomaly detection", GREEN),
    ("3", "35+ Workflows",       "Repayment workflows\nFinancial risk rules\nBehavioral signals",      GOLD),
    ("4", "Butterfly Engine",    "Cascade effect mapping\nChain impact scoring\nAha moment detection", ORANGE),
    ("5", "Score + Reasoning",   "0–100 health score\nGrade A+ to D\nPlain-language AI verdict",      PURPLE),
]

arrow_x_positions = []
for i, (num, title, desc, col) in enumerate(steps):
    cx = Inches(0.4 + i * 2.55)
    cy = Inches(1.7)
    cw = Inches(2.35)
    ch = Inches(4.6)
    # card
    rect(sl, cx, cy, cw, ch,
         fill_color=RGBColor(0x1A,0x2E,0x1A),
         line_color=col, line_width=Pt(1.5))
    # top colour accent
    rect(sl, cx, cy, cw, Inches(0.2), fill_color=col)
    # number circle bg
    rect(sl, cx + Inches(0.85), cy + Inches(0.35), Inches(0.65), Inches(0.65),
         fill_color=col)
    txt(sl, num, cx + Inches(0.85), cy + Inches(0.35), Inches(0.65), Inches(0.65),
        size=20, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    txt(sl, title, cx + Inches(0.12), cy + Inches(1.15), Inches(2.1), Inches(0.55),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc,  cx + Inches(0.12), cy + Inches(1.7),  Inches(2.1), Inches(1.9),
        size=11, color=RGBColor(0x9C,0xA3,0xAF), align=PP_ALIGN.CENTER)
    # arrow (not after last)
    if i < len(steps) - 1:
        arrow_x = cx + cw + Inches(0.02)
        arrow_cy = cy + Inches(2.1)
        txt(sl, "→", arrow_x, arrow_cy, Inches(0.45), Inches(0.5),
            size=22, bold=True, color=RGBColor(0x4B,0x5B,0x4B), align=PP_ALIGN.CENTER)

# Bottom note
txt(sl, "⚡  End-to-end in < 5 seconds  ·  No CIBIL / bureau call required  ·  Works on zero credit history",
    Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5),
    size=12, color=RGBColor(0x6B,0x7A,0x6B), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BUTTERFLY EFFECT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=ORANGE)

section_label(sl, "CORE FEATURE", Inches(0.55), Inches(0.3))
txt(sl, "Butterfly Effect Detection",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=TEXT_DARK)
txt(sl, "How one small event cascades into a portfolio-wide crisis — and how CreditLens catches it early",
    Inches(0.55), Inches(1.45), Inches(12), Inches(0.45),
    size=13, color=MUTED)

divider(sl, Inches(2.0))

# Cascade chain visual
chain = [
    ("Partial\nPayment", ORANGE, "Trigger"),
    ("Penalty\nApplied",  GOLD,   "Impact"),
    ("EMI\nIncreases",    GOLD,   "Impact"),
    ("Next Payment\nMissed", RED,  "Impact"),
    ("Default\nRisk",     RED,    "Outcome"),
]
box_w = Inches(1.9)
box_h = Inches(1.65)
total_w = len(chain) * box_w + (len(chain)-1) * Inches(0.55)
start_x = (W - total_w) / 2

for i, (label, col, stage) in enumerate(chain):
    cx = start_x + i * (box_w + Inches(0.55))
    cy = Inches(2.3)
    card(sl, cx, cy, box_w, box_h, fill=NEAR_WHITE)
    rect(sl, cx, cy, box_w, Inches(0.16), fill_color=col)
    txt(sl, stage, cx, cy - Inches(0.35), box_w, Inches(0.3),
        size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, label, cx, cy + Inches(0.25), box_w, Inches(1.2),
        size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    if i < len(chain)-1:
        ax = cx + box_w + Inches(0.05)
        txt(sl, "→", ax, cy + Inches(0.6), Inches(0.45), Inches(0.5),
            size=22, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# Score impact bar
txt(sl, "Score Impact", Inches(0.55), Inches(4.2), Inches(3), Inches(0.35),
    size=11, bold=True, color=MUTED)
rect(sl, Inches(0.55), Inches(4.6), Inches(12.2), Inches(0.5),
     fill_color=NEAR_WHITE, line_color=BORDER, line_width=Pt(0.5))
rect(sl, Inches(0.55), Inches(4.6), Inches(6.1), Inches(0.5),
     fill_color=RED)
txt(sl, "-32 pts", Inches(0.65), Inches(4.62), Inches(3), Inches(0.35),
    size=13, bold=True, color=WHITE)
txt(sl, "Total cascade score impact from a single missed partial payment",
    Inches(7.0), Inches(4.62), Inches(6), Inches(0.35),
    size=11, color=MUTED)

# Why it matters box
rect(sl, Inches(0.55), Inches(5.3), Inches(12.2), Inches(1.7),
     fill_color=RGBColor(0xFF,0xF7,0xED),
     line_color=ORANGE, line_width=Pt(1))
txt(sl, "Why this matters for lenders",
    Inches(0.8), Inches(5.45), Inches(6), Inches(0.4),
    size=13, bold=True, color=ORANGE)
mtxt(sl, [
    ("Without CreditLens:", True, 12, TEXT_DARK),
    ("  A partner misses one payment → lender notices 90 days later → ₹50L exposure already at risk", False, 11, MUTED),
    ("With CreditLens:", True, 12, GREEN),
    ("  Pattern detected at first partial payment → cascade risk flagged instantly → lender acts in 24 hours", False, 11, MUTED),
], Inches(0.8), Inches(5.85), Inches(12), Inches(1.1), size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DETECTED PATTERNS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=BLUE)

section_label(sl, "PATTERN INTELLIGENCE", Inches(0.55), Inches(0.3))
txt(sl, "What CreditLens Detects",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=TEXT_DARK)
txt(sl, "6 core risk patterns surfaced automatically from transactional data",
    Inches(0.55), Inches(1.45), Inches(12), Inches(0.4),
    size=13, color=MUTED)

divider(sl, Inches(1.95))

patterns = [
    ("⚠️",  "Frequent Partial Payments",  "Borrower consistently pays less than EMI due", "Repayment",  ORANGE),
    ("🔗",  "Multiple Active Loans",       "Debt spread across lenders simultaneously",    "Behavioral", ORANGE),
    ("📉",  "Cash Flow Instability",       "Irregular inflows create repayment risk",       "Financial",  RED),
    ("📈",  "Rising Expenses vs Income",   "Expense growth outpacing revenue",              "Financial",  RED),
    ("🧾",  "GST Filing Irregularity",     "Delayed returns signal revenue volatility",     "Contextual", GOLD),
    ("🌊",  "Cascade Risk",               "One event triggers multiple downstream risks",  "Butterfly",  RED),
]

for i, (icon, name, desc, cat, col) in enumerate(patterns):
    col_idx = i % 3
    row_idx = i // 3
    cx = Inches(0.4  + col_idx * 4.25)
    cy = Inches(2.15 + row_idx * 2.45)
    cw = Inches(4.0)
    ch = Inches(2.2)
    card(sl, cx, cy, cw, ch)
    rect(sl, cx, cy, Inches(0.18), ch, fill_color=col)
    txt(sl, icon, cx + Inches(0.32), cy + Inches(0.2),  Inches(0.6),  Inches(0.55), size=22)
    txt(sl, name, cx + Inches(0.32), cy + Inches(0.75), Inches(3.55), Inches(0.45), size=13, bold=True, color=TEXT_DARK)
    txt(sl, desc, cx + Inches(0.32), cy + Inches(1.2),  Inches(3.55), Inches(0.65), size=11, color=MUTED)
    pill(sl, cat, cx + Inches(0.32), cy + Inches(1.82), Inches(1.2), Inches(0.28),
         RGBColor(0xF3,0xF4,0xF6), MUTED, size=9, bold=False)
    pill(sl, "HIGH RISK" if col==RED else "MODERATE" if col==ORANGE else "MONITOR",
         cx + Inches(1.65), cy + Inches(1.82), Inches(1.3), Inches(0.28),
         RGBColor(0xFF,0xE4,0xE4) if col==RED else RGBColor(0xFF,0xF3,0xE4) if col==ORANGE else RGBColor(0xFE,0xF9,0xC3),
         col, size=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SCORE BREAKDOWN & AI REASONING
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GREEN_DARK)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=GREEN)

section_label(sl, "EXPLAINABLE AI", Inches(0.55), Inches(0.3))
txt(sl, "Every Score is Explainable",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=WHITE)

# Left — score card mock
card(sl, Inches(0.55), Inches(1.75), Inches(5.2), Inches(5.3),
     fill=RGBColor(0x1A,0x2E,0x1A), border=RGBColor(0x22,0x4A,0x22))

txt(sl, "LOAN HEALTH SCORE",
    Inches(0.75), Inches(2.0), Inches(4.5), Inches(0.35),
    size=9, bold=True, color=MUTED)
txt(sl, "72",
    Inches(0.75), Inches(2.35), Inches(2.5), Inches(1.1),
    size=72, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
txt(sl, "/ 100  ·  Grade B+",
    Inches(1.9), Inches(2.88), Inches(3), Inches(0.4),
    size=13, color=RGBColor(0x6B,0x7A,0x6B))

# Score bar
rect(sl, Inches(0.75), Inches(3.55), Inches(4.6), Inches(0.28),
     fill_color=RGBColor(0x0F,0x1D,0x0F), line_color=RGBColor(0x22,0x4A,0x22), line_width=Pt(0.5))
rect(sl, Inches(0.75), Inches(3.55), Inches(3.31), Inches(0.28),
     fill_color=GREEN)

txt(sl, "Proceed with caution — monitor monthly",
    Inches(0.75), Inches(3.95), Inches(4.6), Inches(0.4),
    size=11, color=GOLD)

# Why breakdown bars
breakdown = [
    ("💳 Repayment Impact",   68, BLUE),
    ("🏦 Financial Impact",   45, GREEN),
    ("📊 Behavioural Impact", 80, GOLD),
    ("🌐 Contextual Impact",  30, ORANGE),
]
for j, (lbl, pct, col) in enumerate(breakdown):
    by = Inches(4.5 + j * 0.45)
    txt(sl, lbl, Inches(0.75), by, Inches(2.8), Inches(0.38), size=10, color=RGBColor(0xBB,0xF7,0xD0))
    rect(sl, Inches(3.3), by + Inches(0.08), Inches(1.85), Inches(0.2),
         fill_color=RGBColor(0x0F,0x1D,0x0F))
    rect(sl, Inches(3.3), by + Inches(0.08), Inches(1.85 * pct / 100), Inches(0.2),
         fill_color=col)
    txt(sl, f"{pct}%", Inches(5.2), by, Inches(0.45), Inches(0.38),
        size=10, bold=True, color=col, align=PP_ALIGN.RIGHT)

# Right — AI reasoning mock
card(sl, Inches(6.15), Inches(1.75), Inches(6.85), Inches(5.3),
     fill=RGBColor(0x1A,0x2E,0x1A), border=RGBColor(0x22,0x4A,0x22))

txt(sl, "HOW AI REACHED THIS SCORE",
    Inches(6.35), Inches(1.95), Inches(6.3), Inches(0.35),
    size=9, bold=True, color=MUTED)

wf_steps = [
    ("🔍", "Partial Payment",         "Pattern Detected",  GREEN),
    ("🤖", "GST + Cash Flow cross-check", "AI Analysis",   BLUE),
    ("⚡", "Liquidity stress detected", "Risk Impact",     GOLD),
    ("📊", "-8 Score Impact",          "Score Change",     RED),
]
for j, (icon, step_text, step_label, col) in enumerate(wf_steps):
    sy = Inches(2.45 + j * 1.05)
    # connector line
    if j < len(wf_steps)-1:
        rect(sl, Inches(6.75), sy + Inches(0.45), Pt(2), Inches(0.68),
             fill_color=RGBColor(0x22,0x4A,0x22))
    # node circle
    rect(sl, Inches(6.55), sy + Inches(0.05), Inches(0.42), Inches(0.42),
         fill_color=col)
    txt(sl, icon, Inches(6.55), sy + Inches(0.05), Inches(0.42), Inches(0.42),
        size=13, align=PP_ALIGN.CENTER)
    txt(sl, step_label, Inches(7.1), sy,               Inches(4.5), Inches(0.28),
        size=8, bold=True, color=col)
    txt(sl, step_text,  Inches(7.1), sy + Inches(0.25), Inches(5.5), Inches(0.4),
        size=12, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEMO WALKTHROUGH
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=GREEN)

section_label(sl, "LIVE DEMO", Inches(0.55), Inches(0.3))
txt(sl, "Demo Walkthrough — 5 Steps",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=TEXT_DARK)

divider(sl, Inches(1.55))

demo_steps = [
    ("1", "Open Dashboard",
     "Portfolio overview with 4 KPI cards.\nHealthy loans, at-risk count, avg trust score.",
     GREEN),
    ("2", "Run Partner Health Check",
     "Click 'Run Health Check' on any channel partner.\nWatch the AI pipeline run in real-time.",
     BLUE),
    ("3", "View Trust Score",
     "Score + Grade + Risk Level rendered instantly.\nBusiness verdict: Safe / Caution / Critical.",
     GOLD),
    ("4", "Explore AI Journey Tab",
     "Step-by-step: Pattern → Analysis → Impact → Score.\nClick nodes to expand reasoning.",
     ORANGE),
    ("5", "Cascade Effects Tab",
     "Visual chain: Trigger → Impact → Outcome.\nSee how one event creates portfolio risk.",
     RED),
]

for i, (num, title, desc, col) in enumerate(demo_steps):
    cy = Inches(1.75 + i * 1.06)
    # connector line
    if i < len(demo_steps)-1:
        rect(sl, Inches(0.88), cy + Inches(0.62), Pt(2), Inches(0.55),
             fill_color=BORDER)
    # circle
    rect(sl, Inches(0.55), cy + Inches(0.05), Inches(0.65), Inches(0.65),
         fill_color=col)
    txt(sl, num, Inches(0.55), cy + Inches(0.05), Inches(0.65), Inches(0.65),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.38), cy,               Inches(4.5), Inches(0.38),
        size=13, bold=True, color=TEXT_DARK)
    txt(sl, desc,  Inches(1.38), cy + Inches(0.38), Inches(5.5), Inches(0.6),
        size=11, color=MUTED)

# Right side — "Explain Simply" toggle feature highlight
card(sl, Inches(7.3), Inches(1.75), Inches(5.7), Inches(5.35), fill=NEAR_WHITE)
txt(sl, "💡 Explain Simply — Toggle Mode",
    Inches(7.55), Inches(1.95), Inches(5.2), Inches(0.45),
    size=14, bold=True, color=BLUE)
txt(sl, "Toggle between technical and plain-language explanations in real time:",
    Inches(7.55), Inches(2.45), Inches(5.2), Inches(0.45),
    size=11, color=MUTED)

rect(sl, Inches(7.55), Inches(3.0), Inches(5.0), Inches(1.15),
     fill_color=RGBColor(0xFF,0xF1,0xF1), line_color=RED, line_width=Pt(0.75))
txt(sl, "❌  Technical Mode",
    Inches(7.75), Inches(3.05), Inches(4.5), Inches(0.32),
    size=9, bold=True, color=RED)
txt(sl, '"Liquidity stress due to multi-loan overlap\nand declining DSR ratio"',
    Inches(7.75), Inches(3.37), Inches(4.5), Inches(0.72),
    size=11, italic=True, color=TEXT_DARK)

rect(sl, Inches(7.55), Inches(4.3), Inches(5.0), Inches(1.35),
     fill_color=RGBColor(0xDC,0xFC,0xE7), line_color=GREEN, line_width=Pt(0.75))
txt(sl, "✅  Simple Mode",
    Inches(7.75), Inches(4.35), Inches(4.5), Inches(0.32),
    size=9, bold=True, color=GREEN)
txt(sl, '"Too many loans at once →\nLess money left to repay →\nScore: -8 points"',
    Inches(7.75), Inches(4.67), Inches(4.5), Inches(0.9),
    size=12, bold=True, color=TEXT_DARK)

txt(sl, "Designed for both technical reviewers and business decision-makers",
    Inches(7.55), Inches(5.75), Inches(5.2), Inches(0.45),
    size=10, color=MUTED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MARKET & IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=PURPLE)

section_label(sl, "MARKET OPPORTUNITY", Inches(0.55), Inches(0.3))
txt(sl, "Why Now? Why CreditLens?",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=TEXT_DARK)

divider(sl, Inches(1.55))

# Market numbers
for i, (val, unit, desc, col) in enumerate([
    ("₹37T",  "SAM",  "India MSME formal credit market by 2025",         PURPLE),
    ("₹2.4T", "Gap",  "Annual unmet MSME credit demand",                  RED),
    ("6.3Cr", "MSMEs","Formal MSMEs that need credit but lack bureau data", ORANGE),
]):
    cx = Inches(0.55 + i * 4.25)
    card(sl, cx, Inches(1.75), Inches(3.9), Inches(2.0))
    rect(sl, cx, Inches(1.75), Inches(0.18), Inches(2.0), fill_color=col)
    txt(sl, val,  cx + Inches(0.35), Inches(1.9),  Inches(3.4), Inches(0.8),
        size=36, bold=True, color=col)
    pill(sl, unit, cx + Inches(0.35), Inches(2.7), Inches(0.8), Inches(0.28),
         RGBColor(0xF3,0xF4,0xF6), col, size=10)
    txt(sl, desc, cx + Inches(0.35), Inches(3.05), Inches(3.4), Inches(0.6),
        size=11, color=MUTED)

# Impact metrics
txt(sl, "What CreditLens delivers",
    Inches(0.55), Inches(3.95), Inches(12), Inches(0.4),
    size=14, bold=True, color=TEXT_DARK)

impacts = [
    ("5 sec",   "Full AI score vs 48-hr manual review",    GREEN),
    ("35+",     "Risk workflows running simultaneously",    BLUE),
    ("Zero",    "Bureau/CIBIL dependency",                  ORANGE),
    ("3×",      "Earlier default detection vs traditional", RED),
    ("100%",    "Decision explainability for every score",  PURPLE),
    ("₹50L+",   "Avg exposure protected per flagged partner", GOLD),
]
for i, (val, lbl, col) in enumerate(impacts):
    col_idx = i % 3
    row_idx = i // 3
    cx = Inches(0.55 + col_idx * 4.25)
    cy = Inches(4.45 + row_idx * 1.35)
    card(sl, cx, cy, Inches(4.0), Inches(1.15))
    rect(sl, cx, cy, Inches(0.15), Inches(1.15), fill_color=col)
    txt(sl, val, cx + Inches(0.3), cy + Inches(0.08), Inches(1.5), Inches(0.55),
        size=24, bold=True, color=col)
    txt(sl, lbl, cx + Inches(0.3), cy + Inches(0.62), Inches(3.55), Inches(0.45),
        size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GREEN_DARK)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=BLUE)

section_label(sl, "TECHNOLOGY", Inches(0.55), Inches(0.3))
txt(sl, "What's Under the Hood",
    Inches(0.55), Inches(0.6), Inches(10), Inches(0.8),
    size=36, bold=True, color=WHITE)

layers = [
    ("Frontend",       "Vanilla JS · No framework · Live updates · Responsive", "Zero-dependency UI for maximum speed", BLUE),
    ("Risk Engine",    "35+ modular workflow rules · Category-based scoring", "Repayment · Financial · Behavioral · Contextual", GREEN),
    ("Butterfly AI",   "Cascade chain builder · Trigger-Impact-Outcome mapping", "Detects multi-step default risk propagation", ORANGE),
    ("Pattern Engine", "Real-time signal classification · 10 pattern types", "Derived from transactional + GST + repayment data", GOLD),
    ("Data Layer",     "Digitap bank signals · GST revenue API · Repayment history", "No CIBIL / bureau call — pure alternative data", PURPLE),
]

for i, (layer, desc1, desc2, col) in enumerate(layers):
    cy = Inches(1.75 + i * 1.06)
    rect(sl, Inches(0.55), cy, Inches(12.2), Inches(0.95),
         fill_color=RGBColor(0x1A,0x2E,0x1A),
         line_color=col, line_width=Pt(1.0))
    rect(sl, Inches(0.55), cy, Inches(0.18), Inches(0.95), fill_color=col)
    txt(sl, layer, Inches(0.9), cy + Inches(0.08), Inches(2.0), Inches(0.42),
        size=13, bold=True, color=col)
    txt(sl, desc1, Inches(3.0), cy + Inches(0.08), Inches(5.5), Inches(0.42),
        size=11, color=WHITE)
    txt(sl, desc2, Inches(3.0), cy + Inches(0.52), Inches(5.5), Inches(0.38),
        size=10, color=RGBColor(0x6B,0x7A,0x6B))
    pill(sl, "ACTIVE", Inches(11.5), cy + Inches(0.3), Inches(1.0), Inches(0.3),
         RGBColor(0x16,0xA3,0x4A), WHITE, size=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSE / THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GREEN_DARK)
rect(sl, Inches(0), Inches(0), Inches(0.25), H, fill_color=GREEN)
rect(sl, Inches(0.25), H - Inches(0.25), W - Inches(0.25), Inches(0.25), fill_color=GREEN)

# Large logo text
txt(sl, "CreditLens AI",
    Inches(0.7), Inches(1.5), Inches(12), Inches(1.4),
    size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, Inches(4.5), Inches(3.0), Inches(4.33), Pt(4), fill_color=GREEN)

txt(sl, "Credit Intelligence for the Thin File Economy",
    Inches(0.7), Inches(3.15), Inches(12), Inches(0.6),
    size=20, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.CENTER)

# Three closing points
for i, (icon, line) in enumerate([
    ("⚡", "Zero bureau dependency · works on pure transactional data"),
    ("🌊", "Butterfly effect engine catches cascades before they become defaults"),
    ("🤖", "Every score explained — in technical or plain language"),
]):
    cy = Inches(4.05 + i * 0.7)
    txt(sl, icon, Inches(3.8), cy, Inches(0.5), Inches(0.55), size=18, align=PP_ALIGN.CENTER)
    txt(sl, line, Inches(4.4), cy + Inches(0.05), Inches(5.5), Inches(0.5),
        size=14, color=RGBColor(0x9C,0xA3,0xAF))

txt(sl, "Thank you  ·  Let's demo it live →",
    Inches(0.7), Inches(6.3), Inches(12), Inches(0.65),
    size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ─── SAVE ────────────────────────────────────────────────────────────────────
output = "CreditLens_AI_Demo.pptx"
prs.save(output)
print(f"✅  Saved → {output}  ({len(prs.slides)} slides)")
